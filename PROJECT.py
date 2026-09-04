"""
============================================================
🎓 EFP — EDU FINAL PROJECT
Education & Data Analysis Assistant
============================================================

Provider:
⚡ Groq
✨ Gemini

Fitur:
💬 AI Education Assistant
🎯 Tingkat Penjelasan
📚 Fokus Pembelajaran
📎 Upload File dari Chat

File yang didukung:
📄 TXT
📕 PDF
📝 Word (DOCX)
📊 Excel (XLSX/XLS)
📈 CSV
📽️ PowerPoint (PPTX)
🖼️ PNG/JPG/JPEG

Data Analysis:
📊 Preview Dataset
🔍 Missing Values
📈 Statistik Deskriptif
📉 Visualisasi Data
🧠 AI Data Insight

Cara menjalankan:

streamlit run PROJECT.py
============================================================
"""

import streamlit as st
import pandas as pd

from io import BytesIO
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from langchain_core.messages import (
    AIMessage,
    HumanMessage,
    SystemMessage
)

from langchain_groq import ChatGroq
from langchain_google_genai import ChatGoogleGenerativeAI


# ============================================================
# KONFIGURASI HALAMAN
# ============================================================

st.set_page_config(
    page_title="EFP — Edu Final Project",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded"
)


# ============================================================
# KONSTANTA
# ============================================================

PROVIDER_GROQ = "⚡ Groq"
PROVIDER_GEMINI = "✨ Gemini"


GROQ_MODEL = "openai/gpt-oss-120b"

GEMINI_MODEL = "gemini-2.5-flash"


# ============================================================
# TINGKAT PENJELASAN
# ============================================================

EXPLANATION_LEVELS = [
    "🌱 Pemula",
    "📚 Menengah",
    "🚀 Lanjutan",
    "👨‍🏫 Profesional / Pengajar"
]


# ============================================================
# TOPIK PEMBELAJARAN
# ============================================================

TOPICS = [
    "Umum",
    "Matematika",
    "Bahasa Indonesia",
    "Bahasa Inggris",
    "IPA",
    "IPS",
    "Informatika",
    "Python",
    "Pemrograman",
    "Data Science",
    "Data Analysis",
    "Artificial Intelligence",
    "Machine Learning",
    "Generative AI",
    "Database",
    "Statistika"
]


# ============================================================
# GAYA JAWABAN
# ============================================================

RESPONSE_STYLES = [
    "Santai dan Friendly",
    "Formal dan Profesional",
    "Untuk Pemula",
    "Singkat dan Padat",
    "Detail dan Mendalam"
]


# ============================================================
# CSS
# ============================================================

def load_css():

    st.markdown("""
    <style>

    .main {
        padding-top: 1rem;
    }

    .block-container {
        max-width: 1200px;
        padding-top: 2rem;
        padding-bottom: 2rem;
    }

    .efp-title {
        font-size: 42px;
        font-weight: 800;
        margin-bottom: 0px;
    }

    .efp-subtitle {
        font-size: 18px;
        opacity: 0.7;
        margin-bottom: 20px;
    }

    .feature-card {
        border: 1px solid rgba(128,128,128,0.25);
        border-radius: 15px;
        padding: 20px;
        min-height: 150px;
    }

    .provider-card {
        border: 1px solid rgba(128,128,128,0.25);
        border-radius: 15px;
        padding: 20px;
    }

    </style>
    """, unsafe_allow_html=True)


# ============================================================
# SESSION STATE
# ============================================================

def initialize_session():

    if "app_started" not in st.session_state:
        st.session_state["app_started"] = False


    if "selected_provider" not in st.session_state:
        st.session_state["selected_provider"] = PROVIDER_GROQ


    if "api_keys" not in st.session_state:

        st.session_state["api_keys"] = {
            PROVIDER_GROQ: "",
            PROVIDER_GEMINI: ""
        }


    if "chat_history" not in st.session_state:
        st.session_state["chat_history"] = []


    if "explanation_level" not in st.session_state:
        st.session_state["explanation_level"] = "🌱 Pemula"


    if "topic" not in st.session_state:
        st.session_state["topic"] = "Data Science"


    if "response_style" not in st.session_state:
        st.session_state["response_style"] = "Santai dan Friendly"


    if "temperature" not in st.session_state:
        st.session_state["temperature"] = 0.5


    if "uploaded_dataframe" not in st.session_state:
        st.session_state["uploaded_dataframe"] = None


    if "dataset_name" not in st.session_state:
        st.session_state["dataset_name"] = ""


# ============================================================
# SYSTEM PROMPT
# ============================================================

def create_system_prompt():

    explanation_level = st.session_state["explanation_level"]

    topic = st.session_state["topic"]

    response_style = st.session_state["response_style"]


    return f"""
Kamu adalah EFP (Edu Final Project),
sebuah AI Education & Data Analysis Assistant.

Tugas kamu adalah membantu pengguna untuk:

- Belajar dan memahami materi
- Menjelaskan konsep
- Membuat rangkuman
- Membuat soal dan quiz
- Membantu pemrograman
- Membantu Python
- Membantu Data Science
- Membantu Data Analysis
- Membantu Artificial Intelligence
- Menganalisis dataset
- Memberikan insight dari data
- Membantu memahami dokumen yang diupload

================================================

KONFIGURASI PENGGUNA

🎯 Tingkat Penjelasan:
{explanation_level}

📚 Fokus Pembelajaran:
{topic}

🎭 Gaya Jawaban:
{response_style}

================================================

ATURAN:

1. Selalu jawab menggunakan Bahasa Indonesia.

2. Sesuaikan penjelasan dengan tingkat pengguna.

3. Untuk 🌱 Pemula:
   Gunakan bahasa sederhana dan contoh mudah.

4. Untuk 📚 Menengah:
   Jelaskan konsep dan berikan contoh praktis.

5. Untuk 🚀 Lanjutan:
   Berikan penjelasan teknis dan mendalam.

6. Untuk 👨‍🏫 Profesional / Pengajar:
   Gunakan penjelasan akademis dan profesional.

7. Jika pengguna mengupload dataset:

   - Analisis struktur data
   - Periksa missing values
   - Jelaskan tipe data
   - Analisis statistik
   - Temukan pola
   - Berikan insight
   - Berikan kesimpulan

8. Jangan membuat data atau hasil analisis palsu.

9. Gunakan informasi dari file yang diberikan pengguna.

10. Jika pengguna meminta coding,
    berikan contoh kode yang jelas dan benar.

11. Fokus membantu pengguna belajar,
    memahami materi dan menganalisis data.
"""


# ============================================================
# EXTRACT RESPONSE TEXT
# FIX FORMAT GEMINI
# ============================================================

def extract_response_text(content):

    # Jika sudah string
    if isinstance(content, str):
        return content


    # Jika format list dari Gemini
    if isinstance(content, list):

        result = []

        for item in content:

            if isinstance(item, dict):

                # Format:
                # {"type": "text", "text": "..."}
                if "text" in item:

                    text = item.get("text")

                    if text:
                        result.append(str(text))


            elif isinstance(item, str):

                result.append(item)


        if result:
            return "\n".join(result)


    return str(content)


# ============================================================
# MEMBACA TXT
# ============================================================

def read_txt(file):

    try:

        return file.getvalue().decode(
            "utf-8",
            errors="ignore"
        )

    except Exception as error:

        return f"Gagal membaca TXT: {error}"


# ============================================================
# MEMBACA PDF
# ============================================================

def read_pdf(file):

    try:

        reader = PdfReader(
            BytesIO(file.getvalue())
        )

        text = ""

        for page in reader.pages:

            page_text = page.extract_text()

            if page_text:
                text += page_text + "\n"


        return text


    except Exception as error:

        return f"Gagal membaca PDF: {error}"


# ============================================================
# MEMBACA WORD
# ============================================================

def read_docx(file):

    try:

        document = Document(
            BytesIO(file.getvalue())
        )

        text = ""


        # Paragraph
        for paragraph in document.paragraphs:

            if paragraph.text.strip():

                text += paragraph.text.strip() + "\n"


        # Table
        for table in document.tables:

            for row in table.rows:

                cells = []

                for cell in row.cells:

                    cells.append(
                        cell.text.strip()
                    )

                text += " | ".join(cells) + "\n"


        return text


    except Exception as error:

        return f"Gagal membaca Word: {error}"


# ============================================================
# MEMBACA POWERPOINT
# ============================================================

def read_pptx(file):

    try:

        presentation = Presentation(
            BytesIO(file.getvalue())
        )

        text = ""


        for index, slide in enumerate(
            presentation.slides,
            start=1
        ):

            text += f"\n===== SLIDE {index} =====\n"


            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    if shape.text.strip():

                        text += (
                            shape.text.strip()
                            + "\n"
                        )


        return text


    except Exception as error:

        return f"Gagal membaca PowerPoint: {error}"


# ============================================================
# MEMBACA CSV
# ============================================================

def read_csv(file):

    file_bytes = file.getvalue()


    encodings = [
        "utf-8",
        "utf-8-sig",
        "latin1"
    ]


    for encoding in encodings:

        try:

            dataframe = pd.read_csv(
                BytesIO(file_bytes),
                encoding=encoding
            )

            return dataframe

        except Exception:
            continue


    try:

        dataframe = pd.read_csv(
            BytesIO(file_bytes),
            sep=";"
        )

        return dataframe


    except Exception as error:

        st.error(
            f"❌ Gagal membaca CSV: {error}"
        )

        return None


# ============================================================
# MEMBACA EXCEL
# ============================================================

def read_excel_dataframe(file):

    try:

        dataframe = pd.read_excel(
            BytesIO(file.getvalue())
        )

        return dataframe


    except Exception as error:

        st.error(
            f"❌ Gagal membaca Excel: {error}"
        )

        return None


# ============================================================
# CEK FILE
# ============================================================

def is_image(file):

    filename = file.name.lower()

    return filename.endswith(
        (
            ".png",
            ".jpg",
            ".jpeg"
        )
    )


def is_dataset(file):

    filename = file.name.lower()

    return filename.endswith(
        (
            ".csv",
            ".xlsx",
            ".xls"
        )
    )


def is_document(file):

    filename = file.name.lower()

    return filename.endswith(
        (
            ".txt",
            ".pdf",
            ".docx",
            ".pptx"
        )
    )


# ============================================================
# PROCESS DOCUMENT
# ============================================================

def process_document(file):

    filename = file.name.lower()


    if filename.endswith(".txt"):

        return read_txt(file)


    elif filename.endswith(".pdf"):

        return read_pdf(file)


    elif filename.endswith(".docx"):

        return read_docx(file)


    elif filename.endswith(".pptx"):

        return read_pptx(file)


    return ""


# ============================================================
# DATASET SUMMARY
# ============================================================

def create_dataset_summary(dataframe):

    summary = ""


    summary += """
==============================
INFORMASI DATASET
==============================
"""


    summary += (
        f"Jumlah Baris: {dataframe.shape[0]}\n"
    )

    summary += (
        f"Jumlah Kolom: {dataframe.shape[1]}\n"
    )


    summary += "\nNAMA KOLOM:\n"

    summary += ", ".join(
        dataframe.columns.astype(str)
    )


    summary += "\n\nTIPE DATA:\n"

    summary += (
        dataframe.dtypes.astype(str).to_string()
    )


    summary += "\n\nMISSING VALUES:\n"

    summary += (
        dataframe.isnull().sum().to_string()
    )


    summary += "\n\nCONTOH DATA:\n"

    summary += (
        dataframe.head(10).to_string(
            index=False
        )
    )


    numeric_data = dataframe.select_dtypes(
        include="number"
    )


    if not numeric_data.empty:

        summary += "\n\nSTATISTIK DATA NUMERIK:\n"

        summary += (
            numeric_data.describe().to_string()
        )


    return summary


# ============================================================
# MEMBUAT AI CLIENT
# ============================================================

def create_client():

    provider = st.session_state[
        "selected_provider"
    ]

    api_key = st.session_state[
        "api_keys"
    ][provider]

    temperature = st.session_state[
        "temperature"
    ]


    if provider == PROVIDER_GROQ:

        return ChatGroq(

            model=GROQ_MODEL,

            api_key=api_key,

            temperature=temperature

        )


    elif provider == PROVIDER_GEMINI:

        return ChatGoogleGenerativeAI(

            model=GEMINI_MODEL,

            google_api_key=api_key,

            temperature=temperature

        )


# ============================================================
# LANDING PAGE
# ============================================================

def show_landing_page():

    # --------------------------------------------------------
    # HEADER
    # --------------------------------------------------------

    st.markdown(
        '<div class="efp-title">🎓 EFP</div>',
        unsafe_allow_html=True
    )

    st.markdown(
        '<div class="efp-subtitle">'
        'Edu Final Project — Education & Data Analysis Assistant'
        '</div>',
        unsafe_allow_html=True
    )


    st.write(
        """
Belajar, memahami materi, dan menganalisis data
menjadi lebih mudah dengan bantuan Artificial Intelligence.
"""
    )


    st.divider()


    # --------------------------------------------------------
    # PROVIDER
    # --------------------------------------------------------

    st.subheader("🤖 Pilih AI Provider")

    st.caption(
        "Pilih provider AI yang ingin digunakan."
    )


    col1, col2 = st.columns(2)


    # GROQ
    with col1:

        with st.container(border=True):

            st.subheader("⚡ Groq")

            st.caption(
                "Cepat untuk belajar, diskusi, dan coding."
            )

            st.write("⚡ Respons cepat")

            st.write("💻 Cocok untuk coding")

            st.write("📚 Membantu pembelajaran")

            st.write("📊 Analisis teks dan data")


            if st.button(
                "Pilih ⚡ Groq",
                use_container_width=True
            ):

                st.session_state[
                    "selected_provider"
                ] = PROVIDER_GROQ

                st.rerun()


    # GEMINI
    with col2:

        with st.container(border=True):

            st.subheader("✨ Gemini")

            st.caption(
                "AI fleksibel untuk pembelajaran multimodal."
            )

            st.write("📝 Analisis teks")

            st.write("📄 Membantu memahami dokumen")

            st.write("🖼️ Mendukung analisis gambar")

            st.write("📚 Cocok untuk berbagai materi")


            if st.button(
                "Pilih ✨ Gemini",
                use_container_width=True
            ):

                st.session_state[
                    "selected_provider"
                ] = PROVIDER_GEMINI

                st.rerun()


    st.write("")


    # --------------------------------------------------------
    # API KEY
    # --------------------------------------------------------

    provider = st.session_state[
        "selected_provider"
    ]


    with st.container(border=True):

        st.subheader("🔑 Hubungkan AI")


        if provider == PROVIDER_GROQ:

            st.info(
                f"Provider: {PROVIDER_GROQ}"
            )

            st.caption(
                f"Model: {GROQ_MODEL}"
            )

            api_label = "Groq API Key"

            api_placeholder = (
                "Masukkan Groq API Key..."
            )


        else:

            st.info(
                f"Provider: {PROVIDER_GEMINI}"
            )

            st.caption(
                f"Model: {GEMINI_MODEL}"
            )

            api_label = "Gemini API Key"

            api_placeholder = (
                "Masukkan Gemini API Key..."
            )


        api_key = st.text_input(

            api_label,

            type="password",

            placeholder=api_placeholder

        )


        if st.button(

            "🚀 Mulai EFP",

            type="primary",

            use_container_width=True

        ):

            if api_key:

                st.session_state[
                    "api_keys"
                ][provider] = api_key


                st.session_state[
                    "app_started"
                ] = True


                st.rerun()


            else:

                st.warning(
                    "⚠️ Masukkan API Key terlebih dahulu."
                )


    # --------------------------------------------------------
    # FITUR
    # --------------------------------------------------------

    st.divider()

    st.subheader("🚀 Fitur Utama")


    col1, col2, col3, col4 = st.columns(4)


    with col1:

        with st.container(border=True):

            st.write("### 🎓")

            st.write("**Education AI**")

            st.caption(
                "Belajar dan memahami berbagai materi."
            )


    with col2:

        with st.container(border=True):

            st.write("### 📊")

            st.write("**Data Analysis**")

            st.caption(
                "Analisis dataset CSV dan Excel."
            )


    with col3:

        with st.container(border=True):

            st.write("### 📎")

            st.write("**Upload Materi**")

            st.caption(
                "Upload dokumen langsung dari chat."
            )


    with col4:

        with st.container(border=True):

            st.write("### 🧠")

            st.write("**AI Insight**")

            st.caption(
                "Dapatkan penjelasan dan insight dari AI."
            )


    st.divider()


    st.caption(
        "📄 TXT • 📕 PDF • 📝 Word • 📊 Excel • 📈 CSV • "
        "📽️ PowerPoint • 🖼️ PNG/JPG"
    )


    st.caption(
        "TUGAS FINAL PROJECT : SIGIT PEBRI PRATAMAH"
    )


# ============================================================
# SIDEBAR
# ============================================================

def show_sidebar():

    with st.sidebar:

        st.title("🎓 EFP")

        st.caption(
            "Edu Final Project"
        )


        st.divider()


        # ----------------------------------------------------
        # TINGKAT PENJELASAN
        # ----------------------------------------------------

        st.subheader("🎯 Tingkat Penjelasan")


        current_level = st.session_state[
            "explanation_level"
        ]


        explanation_level = st.selectbox(

            "Pilih Tingkat",

            EXPLANATION_LEVELS,

            index=EXPLANATION_LEVELS.index(
                current_level
            )

        )


        st.session_state[
            "explanation_level"
        ] = explanation_level


        # ----------------------------------------------------
        # TOPIK
        # ----------------------------------------------------

        st.subheader("📚 Fokus Pembelajaran")


        current_topic = st.session_state[
            "topic"
        ]


        topic = st.selectbox(

            "Pilih Topik",

            TOPICS,

            index=TOPICS.index(
                current_topic
            )

        )


        st.session_state["topic"] = topic


        # ----------------------------------------------------
        # GAYA JAWABAN
        # ----------------------------------------------------

        st.subheader("🎭 Gaya Jawaban")


        current_style = st.session_state[
            "response_style"
        ]


        response_style = st.selectbox(

            "Pilih Gaya",

            RESPONSE_STYLES,

            index=RESPONSE_STYLES.index(
                current_style
            )

        )


        st.session_state[
            "response_style"
        ] = response_style


        # ----------------------------------------------------
        # TEMPERATURE
        # ----------------------------------------------------

        st.subheader("🌡️ Kreativitas AI")


        temperature = st.slider(

            "Temperature",

            min_value=0.0,

            max_value=1.0,

            value=st.session_state[
                "temperature"
            ],

            step=0.1,

            help=(
                "Semakin tinggi nilainya, "
                "jawaban AI semakin kreatif."
            )

        )


        st.session_state[
            "temperature"
        ] = temperature


        st.divider()


        # ----------------------------------------------------
        # PROVIDER INFO
        # ----------------------------------------------------

        provider = st.session_state[
            "selected_provider"
        ]


        if provider == PROVIDER_GROQ:

            model = GROQ_MODEL

        else:

            model = GEMINI_MODEL


        st.success(
            f"🟢 Provider aktif: {provider}"
        )


        st.caption(
            f"Model: {model}"
        )


        # ----------------------------------------------------
        # DATASET INFO
        # ----------------------------------------------------

        if st.session_state[
            "uploaded_dataframe"
        ] is not None:

            st.info(
                f"📊 Dataset aktif:\n\n"
                f"{st.session_state['dataset_name']}"
            )


        st.divider()


        # ----------------------------------------------------
        # BUTTON
        # ----------------------------------------------------

        col1, col2 = st.columns(2)


        with col1:

            if st.button(
                "🗑️ Clear",
                use_container_width=True
            ):

                st.session_state[
                    "chat_history"
                ] = []


                st.rerun()


        with col2:

            if st.button(
                "🔑 Ganti Key",
                use_container_width=True
            ):

                provider = st.session_state[
                    "selected_provider"
                ]


                st.session_state[
                    "api_keys"
                ][provider] = ""


                st.session_state[
                    "chat_history"
                ] = []


                st.session_state[
                    "app_started"
                ] = False


                st.rerun()


        # ----------------------------------------------------
        # HAPUS DATASET
        # ----------------------------------------------------

        if st.session_state[
            "uploaded_dataframe"
        ] is not None:

            if st.button(
                "📊 Hapus Dataset",
                use_container_width=True
            ):

                st.session_state[
                    "uploaded_dataframe"
                ] = None


                st.session_state[
                    "dataset_name"
                ] = ""


                st.rerun()


# ============================================================
# DATASET ANALYSIS
# ============================================================

def show_dataset_analysis():

    dataframe = st.session_state[
        "uploaded_dataframe"
    ]


    if dataframe is None:
        return


    st.divider()


    st.subheader(
        f"📊 Dataset: {st.session_state['dataset_name']}"
    )


    tab1, tab2, tab3, tab4 = st.tabs([

        "👀 Preview",
        "📋 Informasi",
        "📊 Statistik",
        "📈 Visualisasi"

    ])


    # --------------------------------------------------------
    # PREVIEW
    # --------------------------------------------------------

    with tab1:

        st.write("### Preview Dataset")

        st.dataframe(
            dataframe,
            use_container_width=True
        )


    # --------------------------------------------------------
    # INFORMASI
    # --------------------------------------------------------

    with tab2:

        col1, col2 = st.columns(2)


        with col1:

            st.metric(
                "Jumlah Baris",
                dataframe.shape[0]
            )


        with col2:

            st.metric(
                "Jumlah Kolom",
                dataframe.shape[1]
            )


        st.write("### Tipe Data")


        data_types = pd.DataFrame({

            "Kolom":
            dataframe.columns.astype(str),

            "Tipe Data":
            dataframe.dtypes.astype(str)

        })


        st.dataframe(
            data_types,
            use_container_width=True
        )


        st.write("### Missing Values")


        missing_values = pd.DataFrame({

            "Kolom":
            dataframe.columns.astype(str),

            "Missing Values":
            dataframe.isnull().sum().values,

            "Persentase (%)":
            (
                dataframe.isnull()
                .mean()
                .values
                * 100
            ).round(2)

        })


        st.dataframe(
            missing_values,
            use_container_width=True
        )


    # --------------------------------------------------------
    # STATISTIK
    # --------------------------------------------------------

    with tab3:

        numeric_data = dataframe.select_dtypes(
            include="number"
        )


        if not numeric_data.empty:

            st.write(
                "### Statistik Deskriptif"
            )


            st.dataframe(
                numeric_data.describe(),
                use_container_width=True
            )


        else:

            st.info(
                "Dataset tidak memiliki kolom numerik."
            )


    # --------------------------------------------------------
    # VISUALISASI
    # --------------------------------------------------------

    with tab4:

        numeric_columns = (
            dataframe
            .select_dtypes(include="number")
            .columns
            .tolist()
        )


        if numeric_columns:

            selected_column = st.selectbox(

                "Pilih Kolom Numerik",

                numeric_columns,

                key="visualization_column"

            )


            chart_data = dataframe[
                selected_column
            ].dropna()


            st.bar_chart(chart_data)


        else:

            st.info(
                "Tidak ada kolom numerik untuk divisualisasikan."
            )


# ============================================================
# MEMBUAT PESAN FILE
# ============================================================

def create_file_context(uploaded_files):

    file_context = ""


    for file in uploaded_files:


        # ----------------------------------------------------
        # DATASET
        # ----------------------------------------------------

        if is_dataset(file):

            dataframe = None


            if file.name.lower().endswith(".csv"):

                dataframe = read_csv(file)


            else:

                dataframe = read_excel_dataframe(
                    file
                )


            if dataframe is not None:


                st.session_state[
                    "uploaded_dataframe"
                ] = dataframe


                st.session_state[
                    "dataset_name"
                ] = file.name


                summary = create_dataset_summary(
                    dataframe
                )


                file_context += f"""

==================================================
DATASET: {file.name}
==================================================

{summary}

"""


        # ----------------------------------------------------
        # DOKUMEN
        # ----------------------------------------------------

        elif is_document(file):

            content = process_document(file)


            if content:


                # Batasi agar prompt tidak terlalu besar
                if len(content) > 30000:

                    content = (
                        content[:30000]
                        +
                        "\n\n[Konten file dipotong karena terlalu panjang]"
                    )


                file_context += f"""

==================================================
FILE: {file.name}
==================================================

{content}

"""


    return file_context


# ============================================================
# MEMBUAT MESSAGE UNTUK AI
# ============================================================

def build_messages(user_prompt, file_context):

    system_prompt = create_system_prompt()


    messages = [

        SystemMessage(
            content=system_prompt
        )

    ]


    # Ambil chat sebelumnya
    messages.extend(
        st.session_state[
            "chat_history"
        ]
    )


    final_prompt = user_prompt


    if file_context:

        final_prompt += f"""

Berikut adalah isi atau informasi dari file yang
diupload pengguna:

{file_context}

Gunakan informasi tersebut sebagai konteks utama
untuk menjawab pertanyaan pengguna.
"""


    messages.append(

        HumanMessage(
            content=final_prompt
        )

    )


    return messages


# ============================================================
# HALAMAN CHAT
# ============================================================

def show_chat():

    # SIDEBAR
    show_sidebar()


    # --------------------------------------------------------
    # HEADER
    # --------------------------------------------------------

    st.title("🎓 EFP")

    st.caption(
        "Education & Data Analysis Assistant"
    )


    provider = st.session_state[
        "selected_provider"
    ]


    model = (
        GROQ_MODEL
        if provider == PROVIDER_GROQ
        else GEMINI_MODEL
    )


    st.caption(
        f"🟢 {provider} • {model}"
    )


    # --------------------------------------------------------
    # CLIENT
    # --------------------------------------------------------

    try:

        client = create_client()


    except Exception as error:

        st.error(
            f"❌ Gagal membuat AI Client: {error}"
        )

        st.stop()


    # --------------------------------------------------------
    # CHAT HISTORY
    # --------------------------------------------------------

    for message in st.session_state[
        "chat_history"
    ]:


        if isinstance(message, HumanMessage):

            with st.chat_message("user"):

                if isinstance(message.content, str):

                    st.markdown(
                        message.content
                    )


        elif isinstance(message, AIMessage):

            with st.chat_message("assistant"):

                response_text = extract_response_text(
                    message.content
                )

                st.markdown(
                    response_text
                )


    # --------------------------------------------------------
    # DATASET ANALYSIS
    # --------------------------------------------------------

    if st.session_state[
        "uploaded_dataframe"
    ] is not None:

        show_dataset_analysis()


    # --------------------------------------------------------
    # CHAT INPUT + FILE
    # --------------------------------------------------------

    chat_input = st.chat_input(

        "Tanyakan sesuatu atau upload materi...",

        accept_file="multiple",

        file_type=[

            "txt",
            "pdf",
            "docx",

            "xlsx",
            "xls",
            "csv",

            "pptx",

            "png",
            "jpg",
            "jpeg"

        ]

    )


    # Jika tidak ada input
    if not chat_input:

        return


    # --------------------------------------------------------
    # AMBIL TEXT DAN FILE
    # --------------------------------------------------------

    user_prompt = chat_input.text

    uploaded_files = chat_input.files


    # Jika hanya upload file
    if not user_prompt and uploaded_files:

        user_prompt = (
            "Tolong analisis file yang saya upload."
        )


    # Jika tidak ada apa-apa
    if not user_prompt:

        return


    # --------------------------------------------------------
    # TAMPILKAN USER MESSAGE
    # --------------------------------------------------------

    with st.chat_message("user"):

        st.markdown(user_prompt)


        if uploaded_files:

            for file in uploaded_files:


                st.caption(
                    f"📎 {file.name}"
                )


                if is_image(file):

                    st.image(
                        file,
                        width=300
                    )


    # --------------------------------------------------------
    # PROSES FILE
    # --------------------------------------------------------

    file_context = create_file_context(
        uploaded_files
    )


    # --------------------------------------------------------
    # WARNING GAMBAR GROQ
    # --------------------------------------------------------

    image_uploaded = any(
        is_image(file)
        for file in uploaded_files
    )


    if (
        image_uploaded
        and provider == PROVIDER_GROQ
    ):

        with st.chat_message("assistant"):

            st.warning(
                "⚠️ Analisis gambar belum didukung "
                "oleh konfigurasi model Groq ini. "
                "Gunakan Gemini untuk menganalisis gambar."
            )


    # --------------------------------------------------------
    # SIMPAN USER MESSAGE
    # --------------------------------------------------------

    st.session_state[
        "chat_history"
    ].append(

        HumanMessage(
            content=user_prompt
        )

    )


    # --------------------------------------------------------
    # BUILD AI MESSAGE
    # --------------------------------------------------------

    messages = build_messages(
        user_prompt,
        file_context
    )


    # --------------------------------------------------------
    # RESPONSE AI
    # --------------------------------------------------------

    with st.chat_message("assistant"):

        with st.spinner(
            "🤔 EFP sedang berpikir..."
        ):

            try:

                response = client.invoke(
                    messages
                )


                response_text = extract_response_text(
                    response.content
                )


                st.markdown(
                    response_text
                )


                # Simpan jawaban
                st.session_state[
                    "chat_history"
                ].append(

                    AIMessage(
                        content=response_text
                    )

                )


            except Exception as error:

                st.error(
                    f"❌ Terjadi error: {error}"
                )


# ============================================================
# MAIN
# ============================================================

initialize_session()

load_css()


# ------------------------------------------------------------
# LANDING PAGE
# ------------------------------------------------------------

if not st.session_state["app_started"]:

    show_landing_page()

    st.stop()


# ------------------------------------------------------------
# CHAT PAGE
# ------------------------------------------------------------

show_chat()
